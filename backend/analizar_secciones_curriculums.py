from __future__ import annotations

import json
import re
import unicodedata
import xml.etree.ElementTree as ET
import zipfile
from collections import Counter
from pathlib import Path
from tempfile import TemporaryDirectory

from pptx import Presentation


BASE_DIR = Path(__file__).resolve().parent

ZIP_PATH = BASE_DIR / "data" / "curriculums_raw" / "CVs.zip"
REPORT_PATH = BASE_DIR / "reporte_preparacion_curriculums.json"
OUTPUT_PATH = BASE_DIR / "reporte_secciones_curriculums.json"


def normalizar_unicode(texto: str | None) -> str:
    """
    Normaliza el texto extraído desde PowerPoint.

    - Convierte Unicode a NFC.
    - Reemplaza caracteres de salto vertical.
    - Normaliza saltos de línea.
    - Elimina espacios al inicio y al final.
    """
    if not texto:
        return ""

    texto = unicodedata.normalize("NFC", texto)
    texto = texto.replace("\x0b", "\n")
    texto = texto.replace("\r\n", "\n")
    texto = texto.replace("\r", "\n")

    return texto.strip()


def normalizar_comparacion(texto: str | None) -> str:
    """
    Convierte el texto a una forma simplificada para comparaciones.

    Ejemplo:
        "Áreas de Especialización"
        -> "areas de especializacion"
    """
    texto_normalizado = normalizar_unicode(texto)

    if not texto_normalizado:
        return ""

    texto_normalizado = unicodedata.normalize(
        "NFKD",
        texto_normalizado,
    )

    texto_normalizado = "".join(
        caracter
        for caracter in texto_normalizado
        if not unicodedata.combining(caracter)
    )

    texto_normalizado = texto_normalizado.lower()
    texto_normalizado = re.sub(
        r"[^a-z0-9]+",
        " ",
        texto_normalizado,
    )

    return " ".join(texto_normalizado.split())


def obtener_textos_desde_xml_pptx(
    ruta_pptx: Path,
) -> list[str]:
    """
    Extrae el texto directamente desde los XML internos del PPTX.

    Esta función sirve como respaldo cuando python-pptx no puede
    abrir el documento por un recurso multimedia corrupto, como
    una imagen con error CRC.

    Solo lee archivos:
        ppt/slides/slideN.xml
    """
    textos: list[str] = []

    namespace = {
        "a": (
            "http://schemas.openxmlformats.org/"
            "drawingml/2006/main"
        )
    }

    with zipfile.ZipFile(ruta_pptx) as archivo_pptx:
        nombres_slides = [
            nombre
            for nombre in archivo_pptx.namelist()
            if re.fullmatch(
                r"ppt/slides/slide\d+\.xml",
                nombre,
            )
        ]

        nombres_slides.sort(
            key=lambda nombre: int(
                re.search(
                    r"slide(\d+)\.xml",
                    nombre,
                ).group(1)
            )
        )

        for nombre_slide in nombres_slides:
            try:
                contenido_xml = archivo_pptx.read(
                    nombre_slide
                )

                raiz = ET.fromstring(
                    contenido_xml
                )

                fragmentos = [
                    nodo.text or ""
                    for nodo in raiz.findall(
                        ".//a:t",
                        namespace,
                    )
                ]

                lineas = [
                    fragmento.strip()
                    for fragmento in fragmentos
                    if fragmento.strip()
                ]

                texto_slide = normalizar_unicode(
                    "\n".join(lineas)
                )

                if texto_slide:
                    textos.append(texto_slide)

            except Exception as exc:
                print(
                    "  Advertencia: no se pudo leer "
                    f"{nombre_slide}: {exc}"
                )

    return textos


def obtener_textos_presentacion(
    ruta_pptx: Path,
) -> tuple[list[str], str]:
    """
    Intenta extraer el texto usando python-pptx.

    Si falla, utiliza como respaldo la lectura directa de los XML
    internos del archivo.

    Retorna:
        textos
        metodo_utilizado
    """
    try:
        presentacion = Presentation(ruta_pptx)
        textos: list[str] = []

        for diapositiva in presentacion.slides:
            for forma in diapositiva.shapes:
                if not hasattr(forma, "text"):
                    continue

                contenido = normalizar_unicode(
                    str(forma.text)
                )

                if contenido:
                    textos.append(contenido)

        return textos, "python-pptx"

    except Exception as exc:
        print(
            "  Advertencia: lectura estándar falló."
        )
        print(
            "  Se intentará recuperar el texto "
            f"desde XML. Motivo: {exc}"
        )

        textos_recuperados = (
            obtener_textos_desde_xml_pptx(
                ruta_pptx
            )
        )

        if not textos_recuperados:
            raise RuntimeError(
                "La lectura estándar falló y tampoco "
                "se pudo recuperar texto desde XML."
            ) from exc

        return textos_recuperados, "xml-recuperado"


def parece_encabezado(linea: str) -> bool:
    """
    Determina si una línea parece ser un encabezado de sección.
    """
    linea = normalizar_unicode(linea)

    if not linea:
        return False

    if len(linea) > 80:
        return False

    if linea.endswith("."):
        return False

    palabras = linea.split()

    if len(palabras) > 8:
        return False

    normalizada = normalizar_comparacion(linea)

    if not normalizada:
        return False

    palabras_clave = (
        "overview",
        "resumen",
        "perfil",
        "experiencia",
        "experiencias",
        "otras experiencias",
        "experiencia profesional",
        "experiencia laboral",
        "areas de especializacion",
        "areas de conocimiento",
        "metodologias",
        "herramientas",
        "tecnologias",
        "clientes asesorados",
        "clientes",
        "estudios",
        "posgrados",
        "formacion",
        "educacion",
        "certificaciones",
        "idiomas",
        "conocimientos",
        "competencias",
        "skills",
    )

    if any(
        palabra_clave in normalizada
        for palabra_clave in palabras_clave
    ):
        return True

    letras = [
        caracter
        for caracter in linea
        if caracter.isalpha()
    ]

    if not letras:
        return False

    porcentaje_mayusculas = sum(
        1
        for caracter in letras
        if caracter.isupper()
    ) / len(letras)

    return porcentaje_mayusculas >= 0.8


def main() -> None:
    """
    Analiza los encabezados encontrados en los CV seleccionados.

    Este script:
    - no modifica Supabase;
    - no importa información;
    - solo lee el ZIP y genera un reporte JSON.
    """
    if not ZIP_PATH.exists():
        raise FileNotFoundError(
            f"No se encontró el ZIP: {ZIP_PATH}"
        )

    if not REPORT_PATH.exists():
        raise FileNotFoundError(
            "No se encontró el reporte previo: "
            f"{REPORT_PATH}"
        )

    reporte_preparacion = json.loads(
        REPORT_PATH.read_text(
            encoding="utf-8"
        )
    )

    seleccionados = reporte_preparacion.get(
        "seleccionados_para_importar",
        [],
    )

    if not seleccionados:
        raise ValueError(
            "No hay currículums seleccionados "
            "en el reporte de preparación."
        )

    frecuencia_encabezados: Counter[str] = Counter()

    archivos_analizados: list[dict] = []
    errores: list[dict] = []
    recuperados_xml: list[dict] = []

    with TemporaryDirectory() as directorio_temporal:
        ruta_temporal_base = Path(
            directorio_temporal
        )

        with zipfile.ZipFile(ZIP_PATH) as archivo_zip:
            indice_zip = {
                Path(nombre).name: nombre
                for nombre in archivo_zip.namelist()
                if nombre.lower().endswith(".pptx")
            }

            total = len(seleccionados)

            for numero, seleccionado in enumerate(
                seleccionados,
                start=1,
            ):
                nombre_archivo = seleccionado[
                    "archivo_seleccionado"
                ]

                print(
                    f"[{numero}/{total}] "
                    f"Analizando {nombre_archivo}"
                )

                nombre_interno = indice_zip.get(
                    nombre_archivo
                )

                if not nombre_interno:
                    errores.append({
                        "archivo": nombre_archivo,
                        "error": (
                            "No se encontró dentro del ZIP."
                        ),
                    })
                    continue

                ruta_temporal = (
                    ruta_temporal_base
                    / nombre_archivo
                )

                try:
                    ruta_temporal.write_bytes(
                        archivo_zip.read(
                            nombre_interno
                        )
                    )

                    textos, metodo_lectura = (
                        obtener_textos_presentacion(
                            ruta_temporal
                        )
                    )

                except Exception as exc:
                    errores.append({
                        "archivo": nombre_archivo,
                        "error": str(exc),
                    })
                    continue

                if metodo_lectura == "xml-recuperado":
                    recuperados_xml.append({
                        "archivo": nombre_archivo,
                        "motivo": (
                            "El archivo fue leído "
                            "directamente desde XML."
                        ),
                    })

                encabezados_archivo: list[str] = []

                for bloque in textos:
                    for linea in bloque.splitlines():
                        linea_limpia = normalizar_unicode(
                            linea
                        )

                        if not parece_encabezado(
                            linea_limpia
                        ):
                            continue

                        encabezado = (
                            normalizar_comparacion(
                                linea_limpia
                            )
                        )

                        if not encabezado:
                            continue

                        encabezados_archivo.append(
                            encabezado
                        )

                        frecuencia_encabezados[
                            encabezado
                        ] += 1

                archivos_analizados.append({
                    "persona_id": seleccionado[
                        "persona_id"
                    ],
                    "persona_nombre": seleccionado[
                        "persona_nombre"
                    ],
                    "archivo": nombre_archivo,
                    "metodo_lectura": metodo_lectura,
                    "encabezados": sorted(
                        set(encabezados_archivo)
                    ),
                    "cantidad_bloques_texto": len(
                        textos
                    ),
                })

    reporte = {
        "curriculums_seleccionados": len(
            seleccionados
        ),
        "curriculums_analizados": len(
            archivos_analizados
        ),
        "errores": len(errores),
        "recuperados_desde_xml": len(
            recuperados_xml
        ),
        "encabezados_frecuentes": [
            {
                "encabezado": encabezado,
                "apariciones": cantidad,
            }
            for encabezado, cantidad
            in frecuencia_encabezados.most_common()
        ],
        "archivos": archivos_analizados,
        "detalle_errores": errores,
        "detalle_recuperados_xml": (
            recuperados_xml
        ),
    }

    OUTPUT_PATH.write_text(
        json.dumps(
            reporte,
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )

    print()
    print("ANÁLISIS FINALIZADO")

    print(
        "CV seleccionados: "
        f"{reporte['curriculums_seleccionados']}"
    )

    print(
        "CV analizados: "
        f"{reporte['curriculums_analizados']}"
    )

    print(
        "Recuperados desde XML: "
        f"{reporte['recuperados_desde_xml']}"
    )

    print(
        f"Errores: {reporte['errores']}"
    )

    print(
        f"Reporte: {OUTPUT_PATH}"
    )

    print()
    print("ENCABEZADOS MÁS FRECUENTES")

    for item in reporte[
        "encabezados_frecuentes"
    ][:30]:
        print(
            f"{item['apariciones']:>3} | "
            f"{item['encabezado']}"
        )


if __name__ == "__main__":
    main()