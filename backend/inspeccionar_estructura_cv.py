from __future__ import annotations

import json
import zipfile
from pathlib import Path
from tempfile import TemporaryDirectory

from pptx import Presentation


BASE_DIR = Path(__file__).resolve().parent
ZIP_PATH = BASE_DIR / "data" / "curriculums_raw" / "CVs.zip"
REPORT_PATH = BASE_DIR / "reporte_preparacion_curriculums.json"
OUTPUT_DIR = BASE_DIR / "data" / "muestras_texto_cv"

CANTIDAD_MUESTRAS = 5


def extraer_textos_presentacion(ruta_pptx: Path) -> list[dict]:
    presentacion = Presentation(ruta_pptx)
    diapositivas: list[dict] = []

    for numero, diapositiva in enumerate(
        presentacion.slides,
        start=1,
    ):
        textos: list[str] = []

        for forma in diapositiva.shapes:
            if not hasattr(forma, "text"):
                continue

            texto = str(forma.text).strip()

            if texto:
                textos.append(texto)

        diapositivas.append({
            "diapositiva": numero,
            "textos": textos,
        })

    return diapositivas


def main() -> None:
    if not ZIP_PATH.exists():
        raise FileNotFoundError(
            f"No se encontró el ZIP: {ZIP_PATH}"
        )

    if not REPORT_PATH.exists():
        raise FileNotFoundError(
            f"No se encontró el reporte: {REPORT_PATH}"
        )

    reporte = json.loads(
        REPORT_PATH.read_text(encoding="utf-8")
    )

    seleccionados = reporte.get(
        "seleccionados_para_importar",
        [],
    )

    if not seleccionados:
        raise ValueError(
            "El reporte no contiene CV seleccionados."
        )

    OUTPUT_DIR.mkdir(
        parents=True,
        exist_ok=True,
    )

    muestras = seleccionados[:CANTIDAD_MUESTRAS]

    with TemporaryDirectory() as carpeta_temporal:
        carpeta_temporal = Path(carpeta_temporal)

        with zipfile.ZipFile(ZIP_PATH) as archivo_zip:
            nombres_zip = {
                Path(nombre).name: nombre
                for nombre in archivo_zip.namelist()
                if nombre.lower().endswith(".pptx")
            }

            for muestra in muestras:
                nombre_archivo = muestra["archivo_seleccionado"]
                nombre_interno = nombres_zip.get(nombre_archivo)

                if not nombre_interno:
                    print(
                        f"No encontrado dentro del ZIP: {nombre_archivo}"
                    )
                    continue

                ruta_temporal = carpeta_temporal / nombre_archivo

                ruta_temporal.write_bytes(
                    archivo_zip.read(nombre_interno)
                )

                contenido = {
                    "persona_id": muestra["persona_id"],
                    "persona_nombre": muestra["persona_nombre"],
                    "archivo": nombre_archivo,
                    "diapositivas": extraer_textos_presentacion(
                        ruta_temporal
                    ),
                }

                nombre_salida = (
                    Path(nombre_archivo).stem + ".json"
                )

                ruta_salida = OUTPUT_DIR / nombre_salida

                ruta_salida.write_text(
                    json.dumps(
                        contenido,
                        ensure_ascii=False,
                        indent=2,
                    ),
                    encoding="utf-8",
                )

                print(f"Generado: {ruta_salida}")

    print()
    print(
        f"Muestras generadas: {len(muestras)}"
    )
    print(f"Carpeta: {OUTPUT_DIR}")


if __name__ == "__main__":
    main()