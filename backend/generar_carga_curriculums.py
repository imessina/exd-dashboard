from __future__ import annotations

import json
import re
import unicodedata
import xml.etree.ElementTree as ET
import zipfile
from collections import OrderedDict
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import Any

from pptx import Presentation


BASE_DIR = Path(__file__).resolve().parent
ZIP_PATH = BASE_DIR / "data" / "curriculums_raw" / "CVs.zip"
REPORT_PREPARACION_PATH = BASE_DIR / "reporte_preparacion_curriculums.json"
OUTPUT_DATA_PATH = BASE_DIR / "data" / "curriculums_carga_actualizada.json"
OUTPUT_REPORT_PATH = BASE_DIR / "reporte_extraccion_curriculums.json"


SECTION_ALIASES: dict[str, tuple[str, ...]] = {
    "resumen": (
        "overview",
        "resumen",
        "resumen profesional",
        "perfil",
        "perfil profesional",
    ),
    "experiencias": (
        "experiencia seleccionadas",
        "experiencias seleccionadas",
        "experiencia seleccionada",
        "experiencia profesional",
        "experiencias profesionales",
        "experiencia laboral",
        "experiencias laborales",
        "experiencia",
        "experiencias",
        "otras experiencias",
    ),
    "areas": (
        "areas de especializacion",
        "areas de conocimiento",
        "metodologias y areas de conocimiento",
        "metodologias y areas de especializacion",
        "metodologias",
        "competencias",
        "skills",
    ),
    "estudios": (
        "estudios y posgrados",
        "estudios y postgrados",
        "estudios",
        "formacion",
        "formacion academica",
        "educacion",
    ),
    "certificaciones": (
        "certificaciones",
        "formacion y certificaciones",
        "certificacion",
    ),
    "herramientas": (
        "herramientas",
        "tecnologias",
        "tecnologias colaborativas",
        "colaborativas",
        "analisis de datos",
        "otras herramientas",
        "metodologias y otros",
        "conocimientos tecnicos",
    ),
    "clientes": (
        "clientes asesorados",
        "clientes asesorados ntt data",
        "clientes",
    ),
    "idiomas": (
        "idiomas",
        "idioma",
    ),
}

ALL_HEADINGS = {
    alias
    for aliases in SECTION_ALIASES.values()
    for alias in aliases
}

DATE_RANGE_PATTERN = re.compile(
    r"(?:\(?\b(?:0?[1-9]|1[0-2])[/\-.](?:19|20)?\d{2}\b\s*[–—-]\s*"
    r"(?:actualidad|presente|(?:0?[1-9]|1[0-2])[/\-.](?:19|20)?\d{2})\)?|"
    r"\(?\b(?:19|20)\d{2}\s*[–—-]\s*(?:actualidad|presente|(?:19|20)\d{2})\)?)",
    re.IGNORECASE,
)

LINKEDIN_PATTERN = re.compile(
    r"https?://(?:www\.)?linkedin\.com/in/[A-Za-z0-9_\-/%?=&.]+",
    re.IGNORECASE,
)

CERT_KEYWORDS = (
    "certified",
    "certificacion",
    "certificación",
    "certificado",
    "certificada",
    "certificado en",
    "practitioner",
    "associate",
    "professional",
    "specialist",
    "especialist",
    "scrum master",
)

STUDY_KEYWORDS = (
    "universidad",
    "instituto",
    "licenciado",
    "licenciada",
    "ingenieria",
    "ingeniería",
    "diplomado",
    "diplomada",
    "magister",
    "máster",
    "master",
    "postgrado",
    "posgrado",
    "titulado",
    "titulada",
    "carrera",
)

NOISE_LINES = {
    "cv project team",
    "ntt data",
    "foto perfil",
    "herramientas",
    "tecnologias",
    "colaborativas",
    "analisis de datos",
    "otras herramientas",
}


def normalize_unicode(text: str | None) -> str:
    if not text:
        return ""

    text = unicodedata.normalize("NFC", str(text))
    text = text.replace("\x0b", "\n")
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = text.replace("\u00a0", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def normalize_compare(text: str | None) -> str:
    text = normalize_unicode(text)
    text = unicodedata.normalize("NFKD", text)
    text = "".join(ch for ch in text if not unicodedata.combining(ch))
    text = text.lower()
    text = re.sub(r"[^a-z0-9]+", " ", text)
    return " ".join(text.split())


def dedupe_preserve_order(values: list[str]) -> list[str]:
    seen: set[str] = set()
    result: list[str] = []

    for value in values:
        clean = normalize_unicode(value).strip(" •-–—\t")
        if not clean:
            continue

        key = normalize_compare(clean)
        if not key or key in seen:
            continue

        seen.add(key)
        result.append(clean)

    return result


def split_nonempty_lines(text: str) -> list[str]:
    return [
        line.strip(" •\t")
        for line in normalize_unicode(text).splitlines()
        if line.strip(" •\t")
    ]


def detect_heading(line: str) -> str | None:
    normalized = normalize_compare(line)

    for section, aliases in SECTION_ALIASES.items():
        for alias in aliases:
            if normalized == alias:
                return section

            if normalized.startswith(alias + " ") and len(normalized) <= len(alias) + 4:
                return section

    return None


def split_inline_heading(line: str) -> tuple[str | None, str]:
    normalized = normalize_compare(line)

    for section, aliases in SECTION_ALIASES.items():
        for alias in sorted(aliases, key=len, reverse=True):
            if normalized == alias:
                return section, ""

            # Try original text with accent-insensitive prefix matching.
            words = line.split()
            for idx in range(1, min(len(words), 8) + 1):
                prefix = " ".join(words[:idx]).rstrip(":")
                if normalize_compare(prefix) == alias:
                    rest = " ".join(words[idx:]).strip(" :\t")
                    return section, rest

    return None, line


def extract_texts_from_xml_pptx(path: Path) -> list[str]:
    texts: list[str] = []
    namespace = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main"
    }

    with zipfile.ZipFile(path) as pptx_zip:
        slide_names = [
            name
            for name in pptx_zip.namelist()
            if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
        ]
        slide_names.sort(
            key=lambda name: int(re.search(r"slide(\d+)\.xml", name).group(1))
        )

        for slide_name in slide_names:
            xml = pptx_zip.read(slide_name)
            root = ET.fromstring(xml)
            fragments = [node.text or "" for node in root.findall(".//a:t", namespace)]
            slide_text = normalize_unicode("\n".join(x for x in fragments if x.strip()))
            if slide_text:
                texts.append(slide_text)

    return texts


def extract_texts(path: Path) -> tuple[list[str], str]:
    try:
        presentation = Presentation(path)
        texts: list[str] = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue

                value = normalize_unicode(str(shape.text))
                if value:
                    texts.append(value)

        return texts, "python-pptx"

    except Exception as exc:
        texts = extract_texts_from_xml_pptx(path)
        if not texts:
            raise RuntimeError(
                "Falló python-pptx y tampoco se recuperó texto desde XML."
            ) from exc
        return texts, "xml-recuperado"


def build_sections(blocks: list[str]) -> dict[str, list[str]]:
    sections: dict[str, list[str]] = {
        name: [] for name in SECTION_ALIASES
    }

    current: str | None = None

    for block in blocks:
        lines = split_nonempty_lines(block)

        for line in lines:
            heading, remainder = split_inline_heading(line)

            if heading:
                current = heading
                if remainder:
                    sections[current].append(remainder)
                continue

            normalized = normalize_compare(line)
            if normalized in NOISE_LINES:
                continue

            if current:
                sections[current].append(line)

    for key in sections:
        sections[key] = dedupe_preserve_order(sections[key])

    return sections


def detect_summary(blocks: list[str], sections: dict[str, list[str]]) -> str | None:
    if sections["resumen"]:
        summary = "\n".join(sections["resumen"]).strip()
        if len(summary) >= 30:
            return summary

    # Some templates put Overview/Resumen and the full body in the same shape.
    for block in blocks:
        lines = split_nonempty_lines(block)
        for idx, line in enumerate(lines):
            if normalize_compare(line) in SECTION_ALIASES["resumen"]:
                collected: list[str] = []
                for next_line in lines[idx + 1 :]:
                    if detect_heading(next_line):
                        break
                    collected.append(next_line)
                summary = "\n".join(collected).strip()
                if len(summary) >= 30:
                    return summary

    return None


def split_list_items(lines: list[str]) -> list[str]:
    items: list[str] = []

    for line in lines:
        for part in re.split(r"\s*[•▪●]\s*|\s*;\s*", line):
            clean = part.strip(" -–—•\t")
            if clean:
                items.append(clean)

    return dedupe_preserve_order(items)


def classify_formation(lines: list[str]) -> tuple[list[str], list[str]]:
    studies: list[str] = []
    certifications: list[str] = []

    for line in lines:
        normalized = normalize_compare(line)
        if any(normalize_compare(k) in normalized for k in CERT_KEYWORDS):
            certifications.append(line)
        elif any(normalize_compare(k) in normalized for k in STUDY_KEYWORDS):
            studies.append(line)
        else:
            # Unknown formation lines are safer as studies than certifications.
            studies.append(line)

    return dedupe_preserve_order(studies), dedupe_preserve_order(certifications)


def parse_experience_header(line: str) -> dict[str, str | None]:
    raw = normalize_unicode(line)
    period_match = DATE_RANGE_PATTERN.search(raw)
    period = period_match.group(0).strip(" ()") if period_match else None

    without_period = raw
    if period_match:
        without_period = (raw[: period_match.start()] + raw[period_match.end() :]).strip()

    parts = [p.strip(" -–—|") for p in re.split(r"\s*\|\s*", without_period) if p.strip(" -–—|")]

    cliente: str | None = None
    proyecto: str | None = None
    rol: str | None = None
    titulo: str | None = without_period or raw

    if len(parts) >= 3:
        cliente = parts[0]
        proyecto = parts[1]
        rol = " | ".join(parts[2:])
        titulo = " | ".join(parts)
    elif len(parts) == 2:
        cliente = parts[0]
        proyecto = parts[1]
        titulo = " | ".join(parts)
    elif " - " in without_period:
        dash_parts = [p.strip() for p in without_period.split(" - ", 1)]
        if len(dash_parts) == 2:
            cliente, proyecto = dash_parts

    return {
        "titulo": titulo or None,
        "cliente": cliente,
        "proyecto": proyecto,
        "rol": rol,
        "periodo": period,
    }


def looks_like_experience_header(line: str) -> bool:
    normalized = normalize_compare(line)

    if DATE_RANGE_PATTERN.search(line):
        return True

    if " | " in line and len(line) <= 220:
        return True

    if any(token in normalized for token in ("proyecto ", "sector ", "banca ", "telecom ")):
        return len(line) <= 220

    return False


def parse_experiences(lines: list[str]) -> list[dict[str, Any]]:
    experiences: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None

    for line in lines:
        clean = normalize_unicode(line)
        if not clean:
            continue

        if looks_like_experience_header(clean):
            if current:
                current["descripcion"] = normalize_unicode("\n".join(current.pop("_descripcion"))) or None
                experiences.append(current)

            current = parse_experience_header(clean)
            current["_descripcion"] = []
        else:
            if current is None:
                # Preserve isolated experience text rather than losing it.
                current = {
                    "titulo": clean[:180],
                    "cliente": None,
                    "proyecto": None,
                    "rol": None,
                    "periodo": None,
                    "_descripcion": [],
                }
            else:
                current["_descripcion"].append(clean)

    if current:
        current["descripcion"] = normalize_unicode("\n".join(current.pop("_descripcion"))) or None
        experiences.append(current)

    # Remove accidental duplicate experience blocks.
    unique: OrderedDict[str, dict[str, Any]] = OrderedDict()
    for exp in experiences:
        key = normalize_compare(
            " | ".join(
                str(exp.get(k) or "")
                for k in ("titulo", "cliente", "proyecto", "rol", "periodo", "descripcion")
            )
        )
        if key and key not in unique:
            unique[key] = exp

    return list(unique.values())


def extract_linkedin(blocks: list[str]) -> str | None:
    joined = "\n".join(blocks)
    match = LINKEDIN_PATTERN.search(joined)
    return match.group(0).rstrip(".,;)") if match else None



def all_clean_lines(blocks: list[str]) -> list[str]:
    lines: list[str] = []
    for block in blocks:
        lines.extend(split_nonempty_lines(block))
    return dedupe_preserve_order(lines)


def extract_section_segments(
    blocks: list[str],
    target_section: str,
) -> list[str]:
    """
    Recupera contenido que aparece después de un encabezado dentro del mismo
    bloque de texto. Complementa build_sections cuando el orden visual de las
    formas del PowerPoint no coincide con el orden lógico de lectura.
    """
    collected: list[str] = []

    for block in blocks:
        lines = split_nonempty_lines(block)
        active = False

        for line in lines:
            heading, remainder = split_inline_heading(line)

            if heading:
                active = heading == target_section
                if active and remainder:
                    collected.append(remainder)
                continue

            if active:
                collected.append(line)

    return dedupe_preserve_order(collected)


def detect_summary_fallback(blocks: list[str]) -> str | None:
    """
    Algunos CV ponen el resumen antes del rótulo 'Resumen' o distribuyen los
    elementos en un orden visual distinto. Busca un párrafo inicial sustancial
    que no parezca una lista de experiencias ni un bloque técnico.
    """
    candidates: list[str] = []

    for block in blocks[:8]:
        text = normalize_unicode(block)
        if len(text) < 120:
            continue

        lines = split_nonempty_lines(text)
        if not lines:
            continue

        normalized = normalize_compare(text)
        date_count = len(DATE_RANGE_PATTERN.findall(text))
        pipe_count = text.count("|")

        if date_count >= 2 or pipe_count >= 4:
            continue
        if any(token in normalized for token in (
            "herramientas tecnologias colaborativas",
            "clientes asesorados",
            "estudios y posgrados",
        )):
            continue

        # Elimina encabezados y datos personales breves.
        meaningful = [
            line for line in lines
            if not detect_heading(line)
            and normalize_compare(line) not in NOISE_LINES
            and len(line) >= 25
        ]
        if meaningful:
            candidates.append("\n".join(meaningful))

    if not candidates:
        return None

    return max(candidates, key=len)


ROLE_KEYWORDS = (
    "engineer", "ingeniero", "ingeniera", "desarrollador", "desarrolladora",
    "designer", "diseñador", "diseñadora", "consultor", "consultora",
    "arquitecto", "arquitecta", "analyst", "analista", "lead", "manager",
    "developer", "ux", "ui", "cro", "prompt", "product", "project",
)


def looks_like_global_experience(line: str) -> bool:
    normalized = normalize_compare(line)
    if len(line) > 260:
        return False
    if DATE_RANGE_PATTERN.search(line):
        return True
    if " | " in line and any(keyword in normalized for keyword in ROLE_KEYWORDS):
        return True
    return False


def extract_global_experiences(blocks: list[str]) -> list[dict[str, Any]]:
    """
    Recupera experiencias que quedaron fuera de la sección lógica debido al
    orden interno de las formas del PPTX. Solo toma líneas con señales fuertes.
    """
    lines = all_clean_lines(blocks)
    experiences: list[dict[str, Any]] = []

    for line in lines:
        if not looks_like_global_experience(line):
            continue
        exp = parse_experience_header(line)
        exp["descripcion"] = None
        experiences.append(exp)

    unique: OrderedDict[str, dict[str, Any]] = OrderedDict()
    for exp in experiences:
        key = normalize_compare(" | ".join(str(exp.get(k) or "") for k in (
            "titulo", "cliente", "proyecto", "rol", "periodo"
        )))
        if key and key not in unique:
            unique[key] = exp

    return list(unique.values())


def merge_experiences(
    primary: list[dict[str, Any]],
    fallback: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    unique: OrderedDict[str, dict[str, Any]] = OrderedDict()

    for exp in primary + fallback:
        key = normalize_compare(" | ".join(str(exp.get(k) or "") for k in (
            "titulo", "cliente", "proyecto", "rol", "periodo"
        )))
        if key and key not in unique:
            unique[key] = exp

    return list(unique.values())


def process_cv(
    selected: dict[str, Any],
    path: Path,
) -> tuple[dict[str, Any], dict[str, Any]]:
    blocks, read_method = extract_texts(path)
    sections = build_sections(blocks)

    summary = detect_summary(blocks, sections) or detect_summary_fallback(blocks)

    studies_direct = split_list_items(sections["estudios"])
    cert_direct = split_list_items(sections["certificaciones"])
    extra_studies, extra_certs = classify_formation(cert_direct)

    studies = dedupe_preserve_order(studies_direct + extra_studies)
    certifications = dedupe_preserve_order(extra_certs)

    areas = split_list_items(
        dedupe_preserve_order(
            sections["areas"] + extract_section_segments(blocks, "areas")
        )
    )
    tools = split_list_items(sections["herramientas"])
    clients = split_list_items(sections["clientes"])
    languages = split_list_items(sections["idiomas"])
    experience_lines = dedupe_preserve_order(
        sections["experiencias"] + extract_section_segments(blocks, "experiencias")
    )
    experiences = merge_experiences(
        parse_experiences(experience_lines),
        extract_global_experiences(blocks),
    )
    linkedin_url = extract_linkedin(blocks)

    warnings: list[str] = []
    if not summary:
        warnings.append("No se detectó resumen profesional.")
    if not experiences:
        warnings.append("No se detectaron experiencias.")
    if not areas:
        warnings.append("No se detectaron áreas de especialización.")
    if read_method == "xml-recuperado":
        warnings.append("El PPTX fue recuperado mediante lectura XML.")

    record: dict[str, Any] = {
        "persona_id": selected["persona_id"],
        "nombre_persona": selected["persona_nombre"],
        "linkedin_url": linkedin_url,
        "curriculum": {
            "resumen_profesional": summary,
            "areas_especializacion": areas,
            "herramientas_tecnologias": tools,
            "clientes_asesorados": clients,
            "estudios_posgrados": studies,
            "idiomas": languages,
            "certificaciones": certifications,
            "archivo_origen": selected["archivo_seleccionado"],
        },
        "experiencias": experiences,
    }

    detail = {
        "persona_id": selected["persona_id"],
        "persona_nombre": selected["persona_nombre"],
        "archivo": selected["archivo_seleccionado"],
        "metodo_lectura": read_method,
        "resumen_detectado": bool(summary),
        "linkedin_detectado": bool(linkedin_url),
        "areas": len(areas),
        "herramientas": len(tools),
        "clientes": len(clients),
        "estudios": len(studies),
        "idiomas": len(languages),
        "certificaciones": len(certifications),
        "experiencias": len(experiences),
        "advertencias": warnings,
    }

    return record, detail


def main() -> None:
    if not ZIP_PATH.exists():
        raise FileNotFoundError(f"No se encontró el ZIP: {ZIP_PATH}")
    if not REPORT_PREPARACION_PATH.exists():
        raise FileNotFoundError(
            f"No se encontró el reporte de preparación: {REPORT_PREPARACION_PATH}"
        )

    report_preparation = json.loads(
        REPORT_PREPARACION_PATH.read_text(encoding="utf-8")
    )
    selected_items = report_preparation.get("seleccionados_para_importar", [])

    if not selected_items:
        raise ValueError("El reporte no contiene currículums seleccionados.")

    output_records: list[dict[str, Any]] = []
    details: list[dict[str, Any]] = []
    errors: list[dict[str, str]] = []

    OUTPUT_DATA_PATH.parent.mkdir(parents=True, exist_ok=True)

    with TemporaryDirectory() as temp_dir:
        temp_dir_path = Path(temp_dir)

        with zipfile.ZipFile(ZIP_PATH) as source_zip:
            zip_index = {
                Path(name).name: name
                for name in source_zip.namelist()
                if name.lower().endswith(".pptx")
            }

            total = len(selected_items)
            for index, selected in enumerate(selected_items, 1):
                filename = selected["archivo_seleccionado"]
                print(f"[{index}/{total}] Extrayendo {filename}")

                internal_name = zip_index.get(filename)
                if not internal_name:
                    errors.append({
                        "archivo": filename,
                        "error": "No se encontró dentro del ZIP.",
                    })
                    continue

                local_path = temp_dir_path / filename

                try:
                    local_path.write_bytes(source_zip.read(internal_name))
                    record, detail = process_cv(selected, local_path)
                    output_records.append(record)
                    details.append(detail)
                except Exception as exc:
                    errors.append({
                        "archivo": filename,
                        "error": str(exc),
                    })

    OUTPUT_DATA_PATH.write_text(
        json.dumps(output_records, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    report = {
        "curriculums_seleccionados": len(selected_items),
        "curriculums_generados": len(output_records),
        "errores": len(errors),
        "total_experiencias": sum(item["experiencias"] for item in details),
        "con_resumen": sum(bool(item["resumen_detectado"]) for item in details),
        "con_linkedin": sum(bool(item["linkedin_detectado"]) for item in details),
        "recuperados_desde_xml": sum(
            item["metodo_lectura"] == "xml-recuperado" for item in details
        ),
        "con_advertencias": sum(bool(item["advertencias"]) for item in details),
        "detalle": details,
        "detalle_errores": errors,
    }

    OUTPUT_REPORT_PATH.write_text(
        json.dumps(report, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    print()
    print("GENERACIÓN FINALIZADA")
    print(f"CV seleccionados: {report['curriculums_seleccionados']}")
    print(f"CV generados: {report['curriculums_generados']}")
    print(f"Total experiencias: {report['total_experiencias']}")
    print(f"Con resumen: {report['con_resumen']}")
    print(f"Con LinkedIn: {report['con_linkedin']}")
    print(f"Recuperados desde XML: {report['recuperados_desde_xml']}")
    print(f"Con advertencias: {report['con_advertencias']}")
    print(f"Errores: {report['errores']}")
    print(f"Archivo de carga: {OUTPUT_DATA_PATH}")
    print(f"Reporte: {OUTPUT_REPORT_PATH}")


if __name__ == "__main__":
    main()
